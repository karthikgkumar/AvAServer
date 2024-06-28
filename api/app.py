import base64
import json
from typing import List, Optional,Type
from fastapi import FastAPI
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from fastapi.responses import RedirectResponse, Response
from io import BytesIO
from pydantic import BaseModel
import openai
import ast
import requests
from dotenv import load_dotenv
from fastapi import FastAPI, HTTPException
import os
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx import Presentation
load_dotenv()  # This loads the variables from .env
from fastapi.responses import StreamingResponse
from mangum import Mangum
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.colors import HexColor, white
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.enums import TA_JUSTIFY, TA_CENTER, TA_LEFT
from io import BytesIO
import logging
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem, Image
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib.styles import ParagraphStyle
from reportlab.platypus.tableofcontents import TableOfContents
os.environ['MPLCONFIGDIR'] = os.getcwd() + "/configs/"
import matplotlib.pyplot as plt
import json
import random
import re

app=FastAPI()
handler=Mangum(app)


class QuestionRequest(BaseModel):
    role: str
    company: str

@app.post("/generate_questions")
async def generate_questions(request: QuestionRequest):
    role = request.role
    company = request.company

    sys_prompt = """
    You are a helpful assistant that can conduct mock interviews for different job roles at various companies. The user will provide the job role and the company, and you will simulate a mock interview for that role and company using the 'conduct_interview' tool.

    When the user specifies a role and company, use the 'conduct_interview' tool with the 'role' and 'company' arguments set to the provided values. The tool will simulate a mock interview based on the given role and company.

    Respond politely and provide the mock interview simulation output from the tool.
    """

    prompt = f"""Generate a python list of 4 interview questions, enclosed in '[]' with 60 word limit as maximum limit tailored for the role of {role} at {company}. The questions should be relevant to the responsibilities, skills, and qualifications required for the role, as well as the company's culture and values. 
    Present the questions strictly as a Python list, with each question as a separate string element enclosed in double quotes "" within the list."""

    msg = [
        {"role": "system", "content": sys_prompt},
        {"role": "user", "content": prompt},
    ]

    openaiclient = openai.Client(
        api_key=os.getenv('OPENAI_API_KEY'),
        base_url=os.getenv('OPENAI_BASE_URL'),
    )
    resp = openaiclient.chat.completions.create(
        messages=msg,
        temperature=0.6,
        model=os.getenv('OPENAI_MODEL'),
        max_tokens=350,
        stream=False,
    )

    question_list = resp.choices[0].message.content

    start_index = question_list.find('[')
    end_index = question_list.find(']') + 1

    bracket_substring = question_list[start_index:end_index]
    questions = ast.literal_eval(bracket_substring)

    return {"questions": questions}

#For the Interview Tool in AvA
@app.post("/generate_pdf")
async def generate_pdf(
    role: str,
    company: Optional[str] = None,
    questions: List[str] = None,
    answer_list: List[str] = None,
    feedback_list: List[str] = None,
    marks_list: List[int] = None,
):
    if not questions or not answer_list or not feedback_list or not marks_list:
        return {"error": "All input parameters (questions, answer_list, feedback_list, marks_list) are required."}

    max_score = len(questions) * 5
    report_filename = f"{role}_{company or ''}_interview_report.pdf"

    # Create an in-memory buffer for the PDF file
    pdf_buffer = BytesIO()

    doc = SimpleDocTemplate(pdf_buffer, pagesize=letter)
    elements = []
    styles = getSampleStyleSheet()
    elements.append(Paragraph(f"Mock Interview Report for {role} at {company or ''}", styles["Heading1"]))
    elements.append(Spacer(1, 12))
    for i, question in enumerate(questions, start=1):
        elements.append(Paragraph(f"Question {i}: {question}", styles["Heading2"]))
        elements.append(Paragraph(f"Your answer: {answer_list[i-1]}", styles["BodyText"]))
        elements.append(Paragraph(f"Evaluation: {feedback_list[i-1]}", styles["BodyText"]))
        elements.append(Paragraph(f"Marks: {marks_list[i-1]}/5", styles["BodyText"]))
        elements.append(Spacer(1, 12))
    elements.append(Paragraph(f"Total Score: {sum(marks_list)}/{max_score}", styles["Heading2"]))

    # Build the PDF and write it to the buffer
    doc.build(elements)
    pdf_value = pdf_buffer.getvalue()
    pdf_buffer.close()
    headers = {"Content-Disposition": f"attachment; filename={report_filename}"}
    response = Response(content=pdf_value, media_type="application/pdf", headers=headers)
    return response

class EvaluationRequest(BaseModel):
    question: str
    answer: str
    role: str
    company: str

class EvaluationResponse(BaseModel):
    feedback: str
    marks: int

@app.post("/evaluate_answer", response_model=EvaluationResponse)
async def evaluate_answer_endpoint(request: EvaluationRequest):
    result = evaluate_answer(request.question, request.answer, request.role, request.company)
    return EvaluationResponse(**result)

def evaluate_answer(question: str, answer: str, role: str, company: str) -> dict:
    sys_prompt = """
    You are a helpful assistant that can conduct mock interviews for different job roles at various companies. The user will provide the job role and the company, and you will simulate a mock interview for that role and company using the 'conduct_interview' tool.

    When the user specifies a role and company, use the 'conduct_interview' tool with the 'role' and 'company' arguments set to the provided values. The tool will simulate a mock interview based on the given role and company.

    Respond politely and provide the mock interview simulation output from the tool.
    """
    prompt = f"""Evaluate the following answer for the role of {role} at {company} and for the 

    question:{question}:

    Answer: {answer} 

    Provide feedback on the answer's quality, relevance, and appropriateness, and assign a score to '<score>' from 0 to 5, where 0 is the lowest and 5 is the highest and assign feedback to '<feedback>'.
        provide the response strictly in the following JSON format with the score and feedback assigned:
                    {{
                        "score": "<score>",
                        "feedback": "<feedback>"
                    }}
        No tool calls should be done.
        """

    syss = sys_prompt
    prompt_n = prompt
    msg = [
        {"role": "system", "content": syss},
        {"role": "user", "content": prompt_n},
    ]

    openaiclient = openai.Client(
        api_key=os.getenv('OPENAI_API_KEY'),
        base_url=os.getenv('OPENAI_BASE_URL'),
    )
    response = openaiclient.chat.completions.create(
        messages=msg,
        temperature=0.6,
        model=os.getenv('OPENAI_MODEL'),
        max_tokens=350,
        stream=False,
    )

    response = response.choices[0].message.content

    try:
        # Extract the JSON object from the response
        json_start = response.find("{")
        json_end = response.rfind("}") + 1
        json_str = response[json_start:json_end]
        evaluation_json = json.loads(json_str)
        return {
            'feedback': evaluation_json.get('feedback', ''),
            'marks': int(evaluation_json.get('score', 0))
        }
    except (json.JSONDecodeError, ValueError):
        return {'feedback': 'Error: Unable to parse the evaluation response', 'marks': 0}


class PresentationRequest(BaseModel):
    topic: str
    num_slides: int
    theme: Optional[str] = "light"

class PresentationCreator:
    def _generate_content(self, topic: str, slide_number: int, total_slides: int) -> tuple:
        openaiclient = openai.Client(
        api_key=os.getenv('OPENAI_API_KEY'),
        base_url=os.getenv('OPENAI_BASE_URL'),
    )
        
        syss = """
            You are a helpful assistant that can create creative presentation slides or PDFs on various topics. The user may ask you to create a presentation on a specific topic, and you should use the 'create_presentation' tool to generate the slides or PDF.

            To create a presentation, use the 'create_presentation' tool with the following arguments:
            - 'topic': The main topic of the presentation
            - 'format': Either 'slides' or 'pdf'
            - 'num_slides': Number of slides or pages to create (1-10)

            For example, if the user asks for a presentation about "Climate Change" with 5 slides, you should call the 'create_presentation' tool with the topic set to "Climate Change", format set to "slides", and num_slides set to 5.

            The content for each slide or page will be automatically generated based on the topic.

            Be sure to respond politely and let the user know if the presentation was created successfully or if there was an error. You can also offer to explain the content of the presentation or ask if they want any modifications.
            """

        prompt_n = f"""
        Create content for a presentation slide about {topic}.
        This is slide {slide_number} out of {total_slides}.
        Provide a title, three key points for the slide, and a prompt for generating an image related to the content.
        Format the output as:
        Title: [Your title here]
        - [First key point]
        - [Second key point]
        - [Third key point]
        Image prompt: [A descriptive prompt for generating an image related to the slide content]

        """
        msg = [
            {"role": "system", "content": syss},
            {"role": "user", "content": prompt_n},
        ]

        response = openaiclient.chat.completions.create(
            messages=msg,
            temperature=0.6,
            model="accounts/fireworks/models/firefunction-v2-rc",
            max_tokens=350,
            stream=False,
        )
        response = response.choices[0].message.content

        lines = response.strip().split('\n')
        title = lines[0].replace('Title: ', '')
        # content = '\n'.join(lines[1:-1])
        content_lines = lines[1:-1]
        formatted_content = []
        for line in content_lines:
            if line.strip():  # Check if the line is not empty
                formatted_line = f"• {line.strip().lstrip('-').strip()}"
                formatted_content.append(formatted_line)
    
        # Join content lines with null characters
        content = '\0'.join(formatted_content)
        image_prompt = lines[-1].replace('Image prompt: ', '')
        
        return title, content, image_prompt
    def _generate_image(self, prompt: str) -> BytesIO:
        image_api_url=os.getenv('IMAGE_API_URL')
        image_api_key=os.getenv('IMAGE_API_KEY')
        headers_image_api = {"Authorization": f"Bearer {image_api_key}"}
        data = {
            "model": "SG161222/Realistic_Vision_V3.0_VAE",
            "negative_prompt": "",
            "prompt": prompt,
            "width": 800,
            "height": 800,
            "steps": 33,
            "n": 1,
            "seed": 8000,
        }
        image_response = requests.post(
            image_api_url, json=data, headers=headers_image_api
        )
        image_response.raise_for_status()
        image_base64 = image_response.json()["output"]["choices"][0]["image_base64"]
        image_bytes = base64.b64decode(image_base64)
    
        return BytesIO(image_bytes)


    
    def _create_slides(self, topic: str, num_slides: int, theme: str) -> str:
        print(f"Creating slides presentation: topic='{topic}', num_slides={num_slides}, theme='{theme}'")
        prs = Presentation()
        
        # Define color schemes
        if theme.lower() == 'dark':
            background_color = RGBColor(32, 33, 36)  # Dark gray
            title_color = RGBColor(255, 255, 255)  # White
            subtitle_color = RGBColor(189, 193, 198)  # Light gray
            text_color = RGBColor(232, 234, 237)  # Off-white
        elif theme.lower() == 'professional':
            background_color = RGBColor(240, 240, 240)  # Light gray
            title_color = RGBColor(31, 73, 125)  # Dark blue
            subtitle_color = RGBColor(68, 114, 196)  # Medium blue
            text_color = RGBColor(0, 0, 0)  # Black
        elif theme.lower() == 'creative':
            background_color = RGBColor(255, 255, 255)  # White
            title_color = RGBColor(255, 67, 67)  # Coral red
            subtitle_color = RGBColor(255, 159, 28)  # Orange
            text_color = RGBColor(87, 117, 144)  # Steel blue
        elif theme.lower() == 'minimalist':
            background_color = RGBColor(255, 255, 255)  # White
            title_color = RGBColor(0, 0, 0)  # Black
            subtitle_color = RGBColor(128, 128, 128)  # Gray
            text_color = RGBColor(64, 64, 64)  # Dark gray
        elif theme.lower() == 'tech':
            background_color = RGBColor(18, 18, 18)  # Very dark gray
            title_color = RGBColor(0, 255, 255)  # Cyan
            subtitle_color = RGBColor(0, 204, 204)  # Darker cyan
            text_color = RGBColor(204, 204, 204)  # Light gray
        else:  # Default light theme
            background_color = RGBColor(255, 255, 255)  # White
            title_color = RGBColor(0, 82, 154)  # Dark blue
            subtitle_color = RGBColor(128, 128, 128)  # Gray
            text_color = RGBColor(64, 64, 64)  # Dark gray
        
        # Define slide layouts
        title_slide_layout = prs.slide_layouts[0]
        content_slide_layout = prs.slide_layouts[5]  # Blank layout
        
        # Create title slide
        title_slide = prs.slides.add_slide(title_slide_layout)
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        
        # Set background color
        background = title_slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = background_color
        
        title.text = topic.upper()
        subtitle.text = f"Generated by AvA"
        
        # Apply styling to title slide
        title.text_frame.paragraphs[0].font.size = Pt(50)
        title.text_frame.paragraphs[0].font.color.rgb = title_color
        subtitle.text_frame.paragraphs[0].font.size = Pt(32)
        subtitle.text_frame.paragraphs[0].font.color.rgb = subtitle_color
        
        for i in range(num_slides - 1):  # -1 because we already created the title slide
            slide = prs.slides.add_slide(content_slide_layout)
            title_slide_content = slide.shapes.title
            
            # Set background color
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = background_color
            
            title, content, image_prompt = self._generate_content(topic, i + 2, num_slides)
            
            title_slide_content.text = title
            # Apply styling to title
            title_slide_content.text_frame.paragraphs[0].font.size = Pt(50)
            title_slide_content.text_frame.paragraphs[0].font.color.rgb = title_color
            
            # Add content as a text box
            left = Inches(0.5)
            top = Inches(2)
            width = Inches(5)
            height = Inches(3)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            
            # Split content by null characters and add each line as a separate paragraph
            content_lines = content.split('\0')
            for idx, line in enumerate(content_lines):
                p = text_frame.add_paragraph()
                p.text = line+'\n'
                p.font.size = Pt(20)
                p.font.color.rgb = text_color
                p.alignment = PP_ALIGN.LEFT
            
            # Generate and add image
            image_stream = self._generate_image(image_prompt)
            left = Inches(5.5)
            top = Inches(2)
            width = Inches(4)
            height = Inches(5)
            slide.shapes.add_picture(image_stream, left, top, width=width, height=height)
        
        # Save the presentation to a BytesIO object
        pptx_file = BytesIO()
        prs.save(pptx_file)
        pptx_file.seek(0)
        
        return pptx_file


@app.post("/create_presentation")
async def create_presentation(request: PresentationRequest):
    creator = PresentationCreator()
    
    try:
        # Create the presentation in memory
        pptx_file = creator._create_slides(request.topic, request.num_slides, request.theme)
        
        # Return the file as a streaming response
        return StreamingResponse(
            pptx_file,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": f'attachment; filename="{request.topic.replace(" ", "_")}_presentation.pptx"'
            }
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    

class PDFCreationRequest(BaseModel):
    topic: str
    num_pages: int

class PDFCreationResponse(BaseModel):
    message: str

@app.post("/create_pdf")
async def create_pdf_endpoint(request: PDFCreationRequest):
    pdf_creation_tool = PDFCreationTool()
    
    try:
        # Create the PDF in memory
        pdf_buffer = pdf_creation_tool._create_pdf(request.topic, request.num_pages)
        
        # Return the file as a streaming response
        return StreamingResponse(
            pdf_buffer,
            media_type="application/pdf",
            headers={
                "Content-Disposition": f'attachment; filename="{request.topic.replace(" ", "_")}_presentation.pdf"'
            }
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


class PDFCreationTool:

    def _create_pdf(self, topic: str, num_pages: int) -> str:
        logging.info(f"Creating PDF presentation: topic='{topic}', num_pages={num_pages}")
        # Create a BytesIO object to store the PDF
        pdf_buffer = BytesIO()
        # Create the PDF document
        doc = SimpleDocTemplate(pdf_buffer, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []

        # Custom styles (same as before)
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=HexColor('#00529A'),  # Dark blue
            alignment=TA_CENTER,
            spaceAfter=12
        )
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['BodyText'],
            fontSize=18,
            textColor=HexColor('#00529A'),  # Dark blue
            alignment=TA_LEFT,
            spaceAfter=12
        )
        subheading_style = ParagraphStyle(
            'CustomSubheading',
            parent=styles['BodyText'],
            fontSize=12,
            textColor=HexColor('#009B77'),  # Green
            alignment=TA_LEFT,
            spaceAfter=8
        )
        body_style = ParagraphStyle(
            'CustomBody',
            parent=styles['BodyText'],
            fontSize=12,
            textColor=HexColor('#404040'),  # Dark gray
            alignment=TA_JUSTIFY,
            spaceAfter=12
        )
        # Add title
        story.append(Paragraph(topic, title_style))
        story.append(Spacer(1, 12))

        # Generate and add title page image
        title_image = self._generate_image(f"Abstract representation of {topic}")
        story.append(Image(title_image, width=6*inch, height=4*inch))
        # story.append(PageBreak())

        for i in range(num_pages - 1):  # -1 because we already created the title page
            logging.info(f"Creating page {i+2}/{num_pages}")
            title, content = self._generate_content(topic, i + 2, num_pages)
            
            story.append(Paragraph(title, heading_style))
            
            # Process content
            paragraphs = content.split('\n\n')
            for para in paragraphs:
                if para.startswith('- '):
                    # Handle bullet points
                    items = [ListItem(Paragraph(item[2:], body_style)) for item in para.split('\n- ')]
                    story.append(ListFlowable(items, bulletType='bullet', start='circle'))
                elif para.startswith('Subheading:'):
                    # Handle subheadings
                    subheading = para[11:]
                    story.append(Paragraph(subheading, subheading_style))
                else:
                    # Regular paragraph with colored important words
                    colored_para = self._color_important_words(para)
                    story.append(Paragraph(colored_para, body_style))
            
            # Generate and add visualization
            chart_image = self._generate_visualization(topic, title)
            if chart_image:
                story.append(Spacer(1, 12))
                story.append(Image(chart_image, width=6*inch, height=4*inch))
            
            # Generate and add image for the current page
            image_prompt = f"Illustration related to {title} in the context of {topic}"
            page_image = self._generate_image(image_prompt)
            story.append(Spacer(1, 12))
            story.append(Image(page_image, width=5*inch, height=3*inch))
            # story.append(PageBreak())

        def add_page_numbers(canvas, doc):
            page_num = canvas.getPageNumber()
            text = f"Page {page_num}"
            canvas.drawRightString(letter[0]-0.5*inch, 0.5*inch, text)

        doc.build(story, onFirstPage=add_page_numbers, onLaterPages=add_page_numbers)
        # Reset the buffer position to the beginning
        pdf_buffer.seek(0)
        
        return pdf_buffer
    
    def _generate_image(self, prompt: str) -> BytesIO:
        image_api_url=os.getenv('IMAGE_API_URL')
        image_api_key=os.getenv('IMAGE_API_KEY')
        headers_image_api = {"Authorization": f"Bearer {image_api_key}"}
        data = {
            "model": "SG161222/Realistic_Vision_V3.0_VAE",
            "negative_prompt": "",
            "prompt": prompt,
            "width": 800,
            "height": 800,
            "steps": 33,
            "n": 1,
            "seed": 8000,
        }
        image_response = requests.post(
            image_api_url, json=data, headers=headers_image_api
        )
        image_response.raise_for_status()
        image_base64 = image_response.json()["output"]["choices"][0]["image_base64"]
        image_bytes = base64.b64decode(image_base64)
        logging.info("Image generated successfully")
        return BytesIO(image_bytes)

    def _generate_content(self, topic: str, page_num: int, total_pages: int):
        prompt = f"""
        Create detailed and well-structured content for page {page_num} of a {total_pages}-page presentation about {topic}.
        The content should be appropriate for a PDF presentation, including:
        1. A clear and concise title for this page
        2. Detailed information about a subtopic related to {topic}
        3. Key points or facts
        4. If applicable, a brief example or case study

        Format the response as follows:
        Title: [Page Title]
        Content: [Detailed content for the page]

        Use paragraphs instead of bullet points where possible. If bullet points are necessary, use '- ' at the start of each point.
        For subheadings, start the line with 'Subheading: '.
        """

        response = self._get_openai_response(prompt)
        
        lines = response.strip().split('\n')
        title = lines[0].replace('Title: ', '')
        content = '\n'.join(lines[2:]).replace('Content: ', '')
        
        return title, content
    
    def _generate_visualization(self, topic: str, title: str) -> Optional[BytesIO]:
        try:
            # Generate data for visualization
            data_prompt = f"Generate sample data for a chart about '{title}' in the context of {topic}. Provide the data as a Python dictionary with 'x' and 'y' keys, where 'x' is a list of labels and 'y' is a list of corresponding values."
            data_response = self._get_openai_response(data_prompt)
            # Try to parse the JSON response
            try:
                data = json.loads(data_response)
            except json.JSONDecodeError:
                # If JSON parsing fails, use fallback data
                logging.warning(f"Failed to parse JSON data for '{title}'. Using fallback data.")
                data = self._generate_fallback_data(title)

            # Ensure data has the correct structure
            if not isinstance(data, dict) or 'x' not in data or 'y' not in data:
                logging.warning(f"Invalid data structure for '{title}'. Using fallback data.")
                data = self._generate_fallback_data(title)

            # Create the chart
            plt.figure(figsize=(8, 6))
            plt.bar(data['x'], data['y'])
            plt.title(title)
            plt.xlabel('Categories')
            plt.ylabel('Values')
            plt.xticks(rotation=45, ha='right')
            plt.tight_layout()

            # Save the chart to a BytesIO object
            img_buffer = BytesIO()
            plt.savefig(img_buffer, format='png')
            img_buffer.seek(0)
            plt.close()

            return img_buffer
        except Exception as e:
            return None
        
    def _generate_fallback_data(self, title: str) -> dict:
        """Generate fallback data if the API response can't be parsed."""
        return {
            'x': [f'Category {i}' for i in range(1, 6)],
            'y': [random.randint(1, 100) for _ in range(5)]
        }

    def _color_important_words(self, text):
        important_words = ['key', 'crucial', 'significant', 'essential', 'fundamental', 'critical', 'vital', 'important']
        for word in important_words:
            text = re.sub(f'\\b{word}\\b', f'<font color="#FF5733">{word}</font>', text, flags=re.IGNORECASE)
        return text

    def _get_openai_response(self, prompt):
        try:
            openaiclient = openai.Client(
                api_key=os.getenv('OPENAI_API_KEY'),
                base_url=os.getenv('OPENAI_BASE_URL'),
            )
            syss = """
            You are a helpful assistant that can create creative PDF presentations on various topics. 
            The user may ask you to create a presentation on a specific topic with a certain number of pages. 
            To create a PDF, use the 'create_pdf' tool with the 'topic' and 'num_pages' arguments.
            For example, if the user says "Create a 5-page presentation about climate change", 
            you should call the 'create_pdf' tool with topic="Climate Change" and num_pages=5.
            Be sure to respond politely and let the user know if the PDF creation was successful or if there was an error.
            """

            msg = [
            {"role": "system", "content": syss},
            {"role": "user", "content": prompt},
        ]
            response = openaiclient.chat.completions.create(
                messages=[{"role": "user", "content": prompt}],
                temperature=0.6,
                model=os.getenv('OPENAI_MODEL'),
                max_tokens=350,
                stream=False,
            )
            return response.choices[0].message.content
        except Exception as e:
            return f"Error generating content: {e}"
    



# Swagger UI is available at http://localhost:8000/docs
@app.get("/", response_class=RedirectResponse)
def redirect_to_docs():
    return "/docs"