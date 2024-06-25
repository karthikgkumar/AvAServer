import base64
import json
from typing import List, Optional
from fastapi import FastAPI
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from fastapi.responses import RedirectResponse, Response
from io import BytesIO
from pydantic import BaseModel
import openai
import ast
import requests
from dotenv import load_dotenv
from fastapi import FastAPI, HTTPException
import os
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx import Presentation
load_dotenv()  # This loads the variables from .env
from fastapi.responses import StreamingResponse

app=FastAPI()


class QuestionRequest(BaseModel):
    role: str
    company: str

@app.post("/generate_questions")
async def generate_questions(request: QuestionRequest):
    role = request.role
    company = request.company

    sys_prompt = """
    You are a helpful assistant that can conduct mock interviews for different job roles at various companies. The user will provide the job role and the company, and you will simulate a mock interview for that role and company using the 'conduct_interview' tool.

    When the user specifies a role and company, use the 'conduct_interview' tool with the 'role' and 'company' arguments set to the provided values. The tool will simulate a mock interview based on the given role and company.

    Respond politely and provide the mock interview simulation output from the tool.
    """

    prompt = f"""Generate a python list of 4 interview questions, enclosed in '[]' with 60 word limit as maximum limit tailored for the role of {role} at {company}. The questions should be relevant to the responsibilities, skills, and qualifications required for the role, as well as the company's culture and values. 
    Present the questions strictly as a Python list, with each question as a separate string element enclosed in double quotes "" within the list."""

    msg = [
        {"role": "system", "content": sys_prompt},
        {"role": "user", "content": prompt},
    ]

    openaiclient = openai.Client(
        api_key=os.getenv('OPENAI_API_KEY'),
        base_url=os.getenv('OPENAI_BASE_URL'),
    )
    resp = openaiclient.chat.completions.create(
        messages=msg,
        temperature=0.6,
        model=os.getenv('OPENAI_MODEL'),
        max_tokens=350,
        stream=False,
    )

    question_list = resp.choices[0].message.content

    start_index = question_list.find('[')
    end_index = question_list.find(']') + 1

    bracket_substring = question_list[start_index:end_index]
    questions = ast.literal_eval(bracket_substring)

    return {"questions": questions}

#For the Interview Tool in AvA
@app.post("/generate_pdf")
async def generate_pdf(
    role: str,
    company: Optional[str] = None,
    questions: List[str] = None,
    answer_list: List[str] = None,
    feedback_list: List[str] = None,
    marks_list: List[int] = None,
):
    if not questions or not answer_list or not feedback_list or not marks_list:
        return {"error": "All input parameters (questions, answer_list, feedback_list, marks_list) are required."}

    max_score = len(questions) * 5
    report_filename = f"{role}_{company or ''}_interview_report.pdf"

    # Create an in-memory buffer for the PDF file
    pdf_buffer = BytesIO()

    doc = SimpleDocTemplate(pdf_buffer, pagesize=letter)
    elements = []
    styles = getSampleStyleSheet()
    elements.append(Paragraph(f"Mock Interview Report for {role} at {company or ''}", styles["Heading1"]))
    elements.append(Spacer(1, 12))
    for i, question in enumerate(questions, start=1):
        elements.append(Paragraph(f"Question {i}: {question}", styles["Heading2"]))
        elements.append(Paragraph(f"Your answer: {answer_list[i-1]}", styles["BodyText"]))
        elements.append(Paragraph(f"Evaluation: {feedback_list[i-1]}", styles["BodyText"]))
        elements.append(Paragraph(f"Marks: {marks_list[i-1]}/5", styles["BodyText"]))
        elements.append(Spacer(1, 12))
    elements.append(Paragraph(f"Total Score: {sum(marks_list)}/{max_score}", styles["Heading2"]))

    # Build the PDF and write it to the buffer
    doc.build(elements)
    pdf_value = pdf_buffer.getvalue()
    pdf_buffer.close()
    headers = {"Content-Disposition": f"attachment; filename={report_filename}"}
    response = Response(content=pdf_value, media_type="application/pdf", headers=headers)
    return response

class EvaluationRequest(BaseModel):
    question: str
    answer: str
    role: str
    company: str

class EvaluationResponse(BaseModel):
    feedback: str
    marks: int

@app.post("/evaluate_answer", response_model=EvaluationResponse)
async def evaluate_answer_endpoint(request: EvaluationRequest):
    result = evaluate_answer(request.question, request.answer, request.role, request.company)
    return EvaluationResponse(**result)

def evaluate_answer(question: str, answer: str, role: str, company: str) -> dict:
    sys_prompt = """
    You are a helpful assistant that can conduct mock interviews for different job roles at various companies. The user will provide the job role and the company, and you will simulate a mock interview for that role and company using the 'conduct_interview' tool.

    When the user specifies a role and company, use the 'conduct_interview' tool with the 'role' and 'company' arguments set to the provided values. The tool will simulate a mock interview based on the given role and company.

    Respond politely and provide the mock interview simulation output from the tool.
    """
    prompt = f"""Evaluate the following answer for the role of {role} at {company} and for the 

    question:{question}:

    Answer: {answer} 

    Provide feedback on the answer's quality, relevance, and appropriateness, and assign a score to '<score>' from 0 to 5, where 0 is the lowest and 5 is the highest and assign feedback to '<feedback>'.
        provide the response strictly in the following JSON format with the score and feedback assigned:
                    {{
                        "score": "<score>",
                        "feedback": "<feedback>"
                    }}
        No tool calls should be done.
        """

    syss = sys_prompt
    prompt_n = prompt
    msg = [
        {"role": "system", "content": syss},
        {"role": "user", "content": prompt_n},
    ]

    openaiclient = openai.Client(
        api_key=os.getenv('OPENAI_API_KEY'),
        base_url=os.getenv('OPENAI_BASE_URL'),
    )
    response = openaiclient.chat.completions.create(
        messages=msg,
        temperature=0.6,
        model=os.getenv('OPENAI_MODEL'),
        max_tokens=350,
        stream=False,
    )

    response = response.choices[0].message.content

    try:
        # Extract the JSON object from the response
        json_start = response.find("{")
        json_end = response.rfind("}") + 1
        json_str = response[json_start:json_end]
        evaluation_json = json.loads(json_str)
        return {
            'feedback': evaluation_json.get('feedback', ''),
            'marks': int(evaluation_json.get('score', 0))
        }
    except (json.JSONDecodeError, ValueError):
        return {'feedback': 'Error: Unable to parse the evaluation response', 'marks': 0}


class PresentationRequest(BaseModel):
    topic: str
    num_slides: int
    theme: Optional[str] = "light"

class PresentationCreator:
    def _generate_content(self, topic: str, slide_number: int, total_slides: int) -> tuple:
        openaiclient = openai.Client(
        api_key=os.getenv('OPENAI_API_KEY'),
        base_url=os.getenv('OPENAI_BASE_URL'),
    )
        
        syss = """
            You are a helpful assistant that can create creative presentation slides or PDFs on various topics. The user may ask you to create a presentation on a specific topic, and you should use the 'create_presentation' tool to generate the slides or PDF.

            To create a presentation, use the 'create_presentation' tool with the following arguments:
            - 'topic': The main topic of the presentation
            - 'format': Either 'slides' or 'pdf'
            - 'num_slides': Number of slides or pages to create (1-10)

            For example, if the user asks for a presentation about "Climate Change" with 5 slides, you should call the 'create_presentation' tool with the topic set to "Climate Change", format set to "slides", and num_slides set to 5.

            The content for each slide or page will be automatically generated based on the topic.

            Be sure to respond politely and let the user know if the presentation was created successfully or if there was an error. You can also offer to explain the content of the presentation or ask if they want any modifications.
            """

        prompt_n = f"""
        Create content for a presentation slide about {topic}.
        This is slide {slide_number} out of {total_slides}.
        Provide a title, three key points for the slide, and a prompt for generating an image related to the content.
        Format the output as:
        Title: [Your title here]
        - [First key point]
        - [Second key point]
        - [Third key point]
        Image prompt: [A descriptive prompt for generating an image related to the slide content]

        """
        msg = [
            {"role": "system", "content": syss},
            {"role": "user", "content": prompt_n},
        ]

        response = openaiclient.chat.completions.create(
            messages=msg,
            temperature=0.6,
            model="accounts/fireworks/models/firefunction-v2-rc",
            max_tokens=350,
            stream=False,
        )
        response = response.choices[0].message.content

        lines = response.strip().split('\n')
        title = lines[0].replace('Title: ', '')
        # content = '\n'.join(lines[1:-1])
        content_lines = lines[1:-1]
        formatted_content = []
        for line in content_lines:
            if line.strip():  # Check if the line is not empty
                formatted_line = f"• {line.strip().lstrip('-').strip()}"
                formatted_content.append(formatted_line)
    
        # Join content lines with null characters
        content = '\0'.join(formatted_content)
        image_prompt = lines[-1].replace('Image prompt: ', '')
        
        return title, content, image_prompt
    def _generate_image(self, prompt: str) -> BytesIO:
        image_api_url=os.getenv('IMAGE_API_URL')
        image_api_key=os.getenv('IMAGE_API_KEY')
        headers_image_api = {"Authorization": f"Bearer {image_api_key}"}
        data = {
            "model": "SG161222/Realistic_Vision_V3.0_VAE",
            "negative_prompt": "",
            "prompt": prompt,
            "width": 800,
            "height": 800,
            "steps": 33,
            "n": 1,
            "seed": 8000,
        }
        image_response = requests.post(
            image_api_url, json=data, headers=headers_image_api
        )
        image_response.raise_for_status()
        image_base64 = image_response.json()["output"]["choices"][0]["image_base64"]
        image_bytes = base64.b64decode(image_base64)
    
        return BytesIO(image_bytes)


    
    def _create_slides(self, topic: str, num_slides: int, theme: str) -> str:
        print(f"Creating slides presentation: topic='{topic}', num_slides={num_slides}, theme='{theme}'")
        prs = Presentation()
        
        # Define color schemes
        if theme.lower() == 'dark':
            background_color = RGBColor(32, 33, 36)  # Dark gray
            title_color = RGBColor(255, 255, 255)  # White
            subtitle_color = RGBColor(189, 193, 198)  # Light gray
            text_color = RGBColor(232, 234, 237)  # Off-white
        elif theme.lower() == 'professional':
            background_color = RGBColor(240, 240, 240)  # Light gray
            title_color = RGBColor(31, 73, 125)  # Dark blue
            subtitle_color = RGBColor(68, 114, 196)  # Medium blue
            text_color = RGBColor(0, 0, 0)  # Black
        elif theme.lower() == 'creative':
            background_color = RGBColor(255, 255, 255)  # White
            title_color = RGBColor(255, 67, 67)  # Coral red
            subtitle_color = RGBColor(255, 159, 28)  # Orange
            text_color = RGBColor(87, 117, 144)  # Steel blue
        elif theme.lower() == 'minimalist':
            background_color = RGBColor(255, 255, 255)  # White
            title_color = RGBColor(0, 0, 0)  # Black
            subtitle_color = RGBColor(128, 128, 128)  # Gray
            text_color = RGBColor(64, 64, 64)  # Dark gray
        elif theme.lower() == 'tech':
            background_color = RGBColor(18, 18, 18)  # Very dark gray
            title_color = RGBColor(0, 255, 255)  # Cyan
            subtitle_color = RGBColor(0, 204, 204)  # Darker cyan
            text_color = RGBColor(204, 204, 204)  # Light gray
        else:  # Default light theme
            background_color = RGBColor(255, 255, 255)  # White
            title_color = RGBColor(0, 82, 154)  # Dark blue
            subtitle_color = RGBColor(128, 128, 128)  # Gray
            text_color = RGBColor(64, 64, 64)  # Dark gray
        
        # Define slide layouts
        title_slide_layout = prs.slide_layouts[0]
        content_slide_layout = prs.slide_layouts[5]  # Blank layout
        
        # Create title slide
        title_slide = prs.slides.add_slide(title_slide_layout)
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        
        # Set background color
        background = title_slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = background_color
        
        title.text = topic.upper()
        subtitle.text = f"Generated by AvA"
        
        # Apply styling to title slide
        title.text_frame.paragraphs[0].font.size = Pt(50)
        title.text_frame.paragraphs[0].font.color.rgb = title_color
        subtitle.text_frame.paragraphs[0].font.size = Pt(32)
        subtitle.text_frame.paragraphs[0].font.color.rgb = subtitle_color
        
        for i in range(num_slides - 1):  # -1 because we already created the title slide
            slide = prs.slides.add_slide(content_slide_layout)
            title_slide_content = slide.shapes.title
            
            # Set background color
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = background_color
            
            title, content, image_prompt = self._generate_content(topic, i + 2, num_slides)
            
            title_slide_content.text = title
            # Apply styling to title
            title_slide_content.text_frame.paragraphs[0].font.size = Pt(50)
            title_slide_content.text_frame.paragraphs[0].font.color.rgb = title_color
            
            # Add content as a text box
            left = Inches(0.5)
            top = Inches(2)
            width = Inches(5)
            height = Inches(3)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            
            # Split content by null characters and add each line as a separate paragraph
            content_lines = content.split('\0')
            for idx, line in enumerate(content_lines):
                p = text_frame.add_paragraph()
                p.text = line+'\n'
                p.font.size = Pt(20)
                p.font.color.rgb = text_color
                p.alignment = PP_ALIGN.LEFT
            
            # Generate and add image
            image_stream = self._generate_image(image_prompt)
            left = Inches(5.5)
            top = Inches(2)
            width = Inches(4)
            height = Inches(5)
            slide.shapes.add_picture(image_stream, left, top, width=width, height=height)
        
        # Save the presentation to a BytesIO object
        pptx_file = BytesIO()
        prs.save(pptx_file)
        pptx_file.seek(0)
        
        return pptx_file


@app.post("/create_presentation")
async def create_presentation(request: PresentationRequest):
    creator = PresentationCreator()
    
    try:
        # Create the presentation in memory
        pptx_file = creator._create_slides(request.topic, request.num_slides, request.theme)
        
        # Return the file as a streaming response
        return StreamingResponse(
            pptx_file,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": f'attachment; filename="{request.topic.replace(" ", "_")}_presentation.pptx"'
            }
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# Swagger UI is available at http://localhost:8000/docs
@app.get("/", response_class=RedirectResponse)
def redirect_to_docs():
    return "/docs"